"""
replace_watermark.py
--------------------
Removes the NotebookLM watermark from AI-generated PPTX files
and replaces it with your own logo.

Usage:
    python replace_watermark.py input.pptx your_logo.png output.pptx

Requirements:
    pip install python-pptx pillow numpy
"""

import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import numpy as np
import io

# ── CONFIG ────────────────────────────────────────────────────────────────────
# Watermark region in 1376x768 pixel space (NotebookLM bottom-right position)
WM_TOP, WM_BOTTOM = 730, 768
WM_LEFT, WM_RIGHT = 1248, 1376

# Logo placement anchor (bottom-right)
LOGO_RIGHT  = 1370
LOGO_BOTTOM = 765
LOGO_WIDTH  = 155   # px — increase/decrease to resize your logo
# ─────────────────────────────────────────────────────────────────────────────


def make_transparent(logo_path):
    """Auto-remove white background from logo if it's a JPG."""
    logo = Image.open(logo_path).convert("RGBA")
    arr = np.array(logo)
    white_mask = (arr[:,:,0] > 240) & (arr[:,:,1] > 240) & (arr[:,:,2] > 240)
    arr[white_mask, 3] = 0
    result = Image.fromarray(arr)
    buf = io.BytesIO()
    result.save(buf, format="PNG")
    buf.seek(0)
    return buf


def sample_background(arr, bottom, left, right):
    """Sample the background color just above the watermark region."""
    ih = arr.shape[0]
    # Take a strip from the absolute bottom of the slide
    strip = arr[ih-8:ih, left:right, :3]
    return tuple(np.median(strip.reshape(-1, 3), axis=0).astype(int))


def process_image(img_bytes, logo_buf):
    img = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
    arr = np.array(img)
    iw, ih = img.size

    # Scale watermark coords to actual image size
    sx, sy = iw / 1376, ih / 768
    top    = int(WM_TOP    * sy)
    bottom = int(WM_BOTTOM * sy)
    left   = int(WM_LEFT   * sx)
    right  = min(iw, int(WM_RIGHT * sx))

    # Sample local background and paint over watermark
    bg = sample_background(arr, top, left, right)
    arr[top:bottom, left:right, :3] = bg
    arr[top:bottom, left:right,  3] = 255

    # Paste logo (transparent) at bottom-right
    logo_buf.seek(0)
    logo = Image.open(logo_buf).convert("RGBA")
    new_w = int(LOGO_WIDTH * sx)
    new_h = int(new_w * (logo.height / logo.width))
    logo  = logo.resize((new_w, new_h), Image.LANCZOS)

    result = Image.fromarray(arr, 'RGBA')
    x = int(LOGO_RIGHT * sx) - new_w
    y = int(LOGO_BOTTOM * sy) - new_h
    result.paste(logo, (x, y), logo)

    buf = io.BytesIO()
    result.save(buf, format="PNG")
    return buf.getvalue()


def replace_watermark(pptx_path, logo_path, output_path):
    logo_buf = make_transparent(logo_path)
    prs = Presentation(pptx_path)

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            blip = shape._element.find(
                './/{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
            )
            rId  = blip.get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
            )
            part = slide.part.related_part(rId)
            part._blob = process_image(part.blob, logo_buf)
            print(f"  Slide {slide_num}")
    prs.save(output_path)
    print(f"Saved -> {output_path}")

if __name__ == "__main__":
    if len(sys.argv) != 4:
        replace_watermark("Modern_SQL_Server_Masterclass.pptx", "Blue_Text_Transparent_Bkg_(1).jpg", "output.pptx")
        sys.exit(1)
    replace_watermark(sys.argv[1], sys.argv[2], sys.argv[3])